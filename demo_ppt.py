#/usr/bin/env python
#coding=utf8

import http.client
import hashlib
from hashlib import md5
import urllib.request, urllib.parse, urllib.error
import random
import json
import string
import time

from pptx import Presentation

prs = Presentation("/Users/albertzhang/Desktop/ppt/internslides.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)

count = 0
appid = input("Enter your app ID: ")
#'20151113000005349'
secretKey = input("Enter your key: ")
#'osubCEzlGjzvw8qdQc41'


httpClient = None
myurl = '/api/trans/vip/translate'
fromLang = 'en'
#input("Enter your source language (refer to language list): ")
toLang = input("Enter your target language (refer to language list): ")
salt = random.randint(32768, 65536)

count = 0
while count < len(text_runs):
    #print(count)
    q = text_runs[count]
    #print(q)
    sign = appid+q+str(salt)+secretKey
    sign = sign.encode('utf-8')
    m1 = hashlib.md5()
    m1.update(sign)
    sign = m1.hexdigest()
    myurl = myurl+'?appid='+appid+'&q='+urllib.parse.quote(q)+'&from='+fromLang+'&to='+toLang+'&salt='+str(salt)+'&sign='+sign

    try:
        httpClient = http.client.HTTPConnection('api.fanyi.baidu.com')
        httpClient.request('GET', myurl)

        #response是HTTPResponse对象
        response = httpClient.getresponse()

        content = str(response.read())

        #print ("All translation info: ", content[2:-1])

        js_content = json.loads(content[2:-1])

        trans_result = js_content['trans_result']

        print(trans_result[-1]['dst'].encode('ascii').decode('unicode-escape'))
        #"Line " + str(count) +": "+
        #print('\u623f\u5b50')
        #count += 1

    except Exception as e:
        print(e)
    finally:
        if httpClient:
            httpClient.close()

    # time.sleep(1)
    count += 1
